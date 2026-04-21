"""
Program Specifications Loader
==============================
Extracts content from any file type a stakeholder might drop into a folder:
  text/data  — TXT, MD, CSV
  documents  — PDF, DOCX/DOC
  spreadsheets — XLSX, XLS
  presentations — PPTX, PPT
  images     — PNG, JPG, JPEG, GIF, BMP, WEBP, TIFF  (via Ollama vision)
  video/audio — MP4, MOV, AVI, MKV, WEBM, MP3, WAV   (via openai-whisper if installed)
"""
from __future__ import annotations

import base64
import json
from pathlib import Path
from typing import Callable

import requests

TEXT_EXTS   = {".txt", ".md", ".rst"}
CSV_EXTS    = {".csv"}
PDF_EXTS    = {".pdf"}
DOCX_EXTS   = {".docx", ".doc"}
EXCEL_EXTS  = {".xlsx", ".xls"}
PPT_EXTS    = {".pptx", ".ppt"}
IMAGE_EXTS  = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tiff", ".tif"}
VIDEO_EXTS  = {".mp4", ".mov", ".avi", ".mkv", ".webm", ".m4v"}
AUDIO_EXTS  = {".mp3", ".wav", ".m4a", ".ogg", ".flac"}

ALL_SUPPORTED = (
    TEXT_EXTS | CSV_EXTS | PDF_EXTS | DOCX_EXTS | EXCEL_EXTS |
    PPT_EXTS | IMAGE_EXTS | VIDEO_EXTS | AUDIO_EXTS
)


# ---------------------------------------------------------------------------
# Individual extractors
# ---------------------------------------------------------------------------

def _extract_text(path: Path) -> tuple[str, str | None]:
    try:
        return path.read_text(encoding="utf-8", errors="ignore"), None
    except Exception as exc:
        return "", str(exc)


def _extract_csv(path: Path) -> tuple[str, str | None]:
    try:
        import pandas as pd
        df = pd.read_csv(path, nrows=500)
        header = f"[CSV: {path.name} — {len(df)} rows × {len(df.columns)} cols]\n"
        return header + df.to_string(index=False, max_rows=200), None
    except Exception as exc:
        return "", str(exc)


def _extract_pdf(path: Path) -> tuple[str, str | None]:
    try:
        import pdfplumber
    except ImportError:
        return "", "pdfplumber not installed"
    try:
        with pdfplumber.open(str(path)) as pdf:
            text = "\n".join(p.extract_text() or "" for p in pdf.pages[:40])
        return text, None
    except Exception as exc:
        return "", str(exc)


def _extract_docx(path: Path) -> tuple[str, str | None]:
    try:
        import docx
    except ImportError:
        return "", "python-docx not installed"
    try:
        doc = docx.Document(str(path))
        parts = [p.text for p in doc.paragraphs]
        for table in doc.tables:
            for row in table.rows:
                parts.append(" | ".join(cell.text for cell in row.cells))
        return "\n".join(parts), None
    except Exception as exc:
        return "", str(exc)


def _extract_excel(path: Path) -> tuple[str, str | None]:
    try:
        import openpyxl
    except ImportError:
        return "", "openpyxl not installed — run: pip install openpyxl"
    try:
        wb = openpyxl.load_workbook(str(path), data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"[Sheet: {sheet_name}]")
            row_count = 0
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join("" if v is None else str(v) for v in row)
                if row_text.strip():
                    parts.append(row_text)
                    row_count += 1
                    if row_count >= 300:
                        parts.append(f"... (truncated after 300 rows)")
                        break
        return "\n".join(parts), None
    except Exception as exc:
        return "", str(exc)


def _extract_pptx(path: Path) -> tuple[str, str | None]:
    try:
        from pptx import Presentation
    except ImportError:
        return "", "python-pptx not installed — run: pip install python-pptx"
    try:
        prs = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            slide_texts.append(t)
            if slide_texts:
                parts.append(f"[Slide {i}]\n" + "\n".join(slide_texts))
        return "\n\n".join(parts), None
    except Exception as exc:
        return "", str(exc)


def _extract_image_ollama(path: Path, ollama_url: str, model: str) -> tuple[str, str | None]:
    """Send image to Ollama vision model and get a description."""
    try:
        with open(path, "rb") as f:
            b64 = base64.b64encode(f.read()).decode()
        payload = {
            "model": model,
            "messages": [
                {
                    "role": "user",
                    "content": (
                        "You are helping build a curriculum. Describe everything visible in this image "
                        "in detail: any text, tables, diagrams, charts, requirements, learning objectives, "
                        "competency frameworks, or stakeholder notes. Transcribe all readable text verbatim."
                    ),
                    "images": [b64],
                }
            ],
            "stream": False,
            "options": {"num_ctx": 4096},
        }
        resp = requests.post(f"{ollama_url}/api/chat", json=payload, timeout=120)
        resp.raise_for_status()
        content = resp.json().get("message", {}).get("content", "")
        if not content:
            return "", "No description returned from vision model"
        return f"[Image analysis: {path.name}]\n{content}", None
    except Exception as exc:
        return "", f"Vision model error: {exc}"


def _extract_video_whisper(path: Path) -> tuple[str, str | None]:
    """Transcribe video/audio using openai-whisper (if installed)."""
    try:
        import whisper  # type: ignore
    except ImportError:
        return "", "openai-whisper not installed — run: pip install openai-whisper"
    try:
        model = whisper.load_model("base")
        result = model.transcribe(str(path))
        transcript = result.get("text", "").strip()
        if not transcript:
            return "", "Whisper returned empty transcript"
        return f"[Video transcript: {path.name}]\n{transcript}", None
    except Exception as exc:
        return "", f"Whisper transcription error: {exc}"


# ---------------------------------------------------------------------------
# Public interface
# ---------------------------------------------------------------------------

def _file_type_label(path: Path) -> str:
    ext = path.suffix.lower()
    if ext in TEXT_EXTS:        return "text"
    if ext in CSV_EXTS:         return "spreadsheet"
    if ext in PDF_EXTS:         return "pdf"
    if ext in DOCX_EXTS:        return "word"
    if ext in EXCEL_EXTS:       return "excel"
    if ext in PPT_EXTS:         return "presentation"
    if ext in IMAGE_EXTS:       return "image"
    if ext in VIDEO_EXTS:       return "video"
    if ext in AUDIO_EXTS:       return "audio"
    return "unknown"


def load_program_specs(
    folder_path: str,
    ollama_url: str = "http://localhost:11434",
    vision_model: str | None = None,
    progress_callback: Callable[[int, int, str], None] | None = None,
) -> list[dict]:
    """
    Recursively scan folder_path and extract content from every recognised file.

    Returns list of dicts:
      filename, path, file_type, text, char_count, error
    """
    folder = Path(folder_path)
    if not folder.exists():
        return []

    files = sorted(
        p for p in folder.rglob("*")
        if p.is_file() and p.suffix.lower() in ALL_SUPPORTED
    )

    results = []
    for i, path in enumerate(files):
        if progress_callback:
            progress_callback(i, len(files), path.name)

        ext = path.suffix.lower()
        file_type = _file_type_label(path)

        if ext in TEXT_EXTS:
            text, error = _extract_text(path)
        elif ext in CSV_EXTS:
            text, error = _extract_csv(path)
        elif ext in PDF_EXTS:
            text, error = _extract_pdf(path)
        elif ext in DOCX_EXTS:
            text, error = _extract_docx(path)
        elif ext in EXCEL_EXTS:
            text, error = _extract_excel(path)
        elif ext in PPT_EXTS:
            text, error = _extract_pptx(path)
        elif ext in IMAGE_EXTS:
            if vision_model:
                text, error = _extract_image_ollama(path, ollama_url, vision_model)
            else:
                text, error = "", "No vision model selected — choose one to analyse images"
        elif ext in (VIDEO_EXTS | AUDIO_EXTS):
            text, error = _extract_video_whisper(path)
        else:
            text, error = "", f"Unsupported: {ext}"

        results.append({
            "filename": path.name,
            "path": str(path),
            "file_type": file_type,
            "text": text,
            "char_count": len(text),
            "error": error,
        })

    return results
