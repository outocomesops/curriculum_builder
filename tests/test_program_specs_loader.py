from pathlib import Path
from unittest.mock import patch

from loaders import program_specs_loader as psl
from tests.conftest import FakeResponse


def test_file_type_label_known_extensions(tmp_path: Path):
    assert psl._file_type_label(tmp_path / "a.txt") == "text"
    assert psl._file_type_label(tmp_path / "a.csv") == "spreadsheet"
    assert psl._file_type_label(tmp_path / "a.pdf") == "pdf"
    assert psl._file_type_label(tmp_path / "a.docx") == "word"
    assert psl._file_type_label(tmp_path / "a.xlsx") == "excel"
    assert psl._file_type_label(tmp_path / "a.pptx") == "presentation"
    assert psl._file_type_label(tmp_path / "a.png") == "image"
    assert psl._file_type_label(tmp_path / "a.mp4") == "video"
    assert psl._file_type_label(tmp_path / "a.mp3") == "audio"
    assert psl._file_type_label(tmp_path / "a.zzz") == "unknown"


def test_extract_text_reads_utf8(tmp_path: Path):
    f = tmp_path / "note.txt"
    f.write_text("simple note", encoding="utf-8")
    text, err = psl._extract_text(f)
    assert text == "simple note"
    assert err is None


def test_extract_csv_renders_table(tmp_path: Path):
    f = tmp_path / "data.csv"
    f.write_text("a,b\n1,2\n3,4\n", encoding="utf-8")
    text, err = psl._extract_csv(f)
    assert err is None
    assert "[CSV: data.csv" in text
    assert "1" in text and "4" in text


def test_extract_docx_reads_paragraphs_and_tables(tmp_path: Path):
    from docx import Document
    f = tmp_path / "spec.docx"
    doc = Document()
    doc.add_paragraph("Paragraph body.")
    table = doc.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "c1"
    table.rows[0].cells[1].text = "c2"
    doc.save(str(f))
    text, err = psl._extract_docx(f)
    assert err is None
    assert "Paragraph body." in text
    assert "c1" in text and "c2" in text


def test_extract_excel_reads_sheets(tmp_path: Path):
    import openpyxl
    f = tmp_path / "sheet.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Summary"
    ws.append(["Name", "Value"])
    ws.append(["metric_a", 7])
    wb.save(str(f))
    text, err = psl._extract_excel(f)
    assert err is None
    assert "[Sheet: Summary]" in text
    assert "metric_a" in text


def test_extract_pptx_reads_slides(tmp_path: Path):
    from pptx import Presentation
    from pptx.util import Inches

    f = tmp_path / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tx_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tx_box.text_frame.text = "Slide heading"
    prs.save(str(f))
    text, err = psl._extract_pptx(f)
    assert err is None
    assert "Slide 1" in text
    assert "Slide heading" in text


def test_load_program_specs_missing_folder(tmp_path: Path):
    assert psl.load_program_specs(str(tmp_path / "nope")) == []


def test_load_program_specs_mixed_files(tmp_path: Path):
    folder = tmp_path / "specs"
    folder.mkdir()
    (folder / "a.txt").write_text("alpha", encoding="utf-8")
    (folder / "b.csv").write_text("x,y\n1,2\n", encoding="utf-8")
    (folder / "skip.xyz").write_text("ignored", encoding="utf-8")
    # image file with no vision model should return an error
    (folder / "pic.png").write_bytes(b"\x89PNG\r\n")

    calls = []

    def cb(i, total, name):
        calls.append((i, total, name))

    results = psl.load_program_specs(str(folder), progress_callback=cb)
    by_name = {r["filename"]: r for r in results}
    assert "a.txt" in by_name and by_name["a.txt"]["error"] is None
    assert "b.csv" in by_name
    assert "skip.xyz" not in by_name  # unsupported and skipped entirely
    assert "pic.png" in by_name and by_name["pic.png"]["error"]
    assert len(calls) == len(results)


@patch("loaders.program_specs_loader.requests.post")
def test_extract_image_ollama_success(mock_post, tmp_path: Path):
    mock_post.return_value = FakeResponse(
        json_data={"message": {"content": "A chart."}}, status_code=200
    )
    img = tmp_path / "pic.png"
    img.write_bytes(b"fake-png")
    text, err = psl._extract_image_ollama(img, "http://localhost:11434", "llava")
    assert err is None
    assert "A chart." in text


@patch("loaders.program_specs_loader.requests.post")
def test_extract_image_ollama_empty_response(mock_post, tmp_path: Path):
    mock_post.return_value = FakeResponse(json_data={"message": {"content": ""}}, status_code=200)
    img = tmp_path / "pic.png"
    img.write_bytes(b"fake")
    text, err = psl._extract_image_ollama(img, "http://localhost:11434", "llava")
    assert text == ""
    assert err is not None


@patch("loaders.program_specs_loader.requests.post")
def test_extract_image_ollama_http_error(mock_post, tmp_path: Path):
    mock_post.side_effect = RuntimeError("connection refused")
    img = tmp_path / "pic.png"
    img.write_bytes(b"fake")
    text, err = psl._extract_image_ollama(img, "http://localhost:11434", "llava")
    assert text == ""
    assert "Vision model error" in err
